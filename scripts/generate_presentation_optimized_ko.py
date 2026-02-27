#!/usr/bin/env python3
"""Generate an optimized Korean presentation (clear 기승전결, reduced slide count)."""

from __future__ import annotations

import csv
import json
from pathlib import Path
from typing import Dict, List, Sequence

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas


ROOT = Path(".")
REPORTS = ROOT / "derived" / "reports"

MANIFEST_SUMMARY = ROOT / "derived" / "manifests" / "summary.json"
PHASE2_GLOBAL = REPORTS / "phase2_global_metrics.csv"
PHASE2_BOOTSTRAP = REPORTS / "phase2_pairwise_bootstrap.csv"
PHASE3_GLOBAL = REPORTS / "phase3_global_metrics.csv"
PHASE35_V5_MAIN = REPORTS / "phase35_strong_v1_main_metrics.csv"
PHASE35_V7 = REPORTS / "phase35_next_v7_metrics.csv"
PHASE35_V8 = REPORTS / "phase35_next_v8_metrics.csv"
PHASE35_CROSS = REPORTS / "phase35_cross_domain_adapt_metrics.csv"

OUT_PPTX = REPORTS / "research_progress_summary_2026-02-24_optimized_ko.pptx"
OUT_PDF = REPORTS / "research_progress_summary_2026-02-24_optimized_ko.pdf"

FONT_NAME = "Malgun Gothic"
TITLE_COLOR = RGBColor(17, 38, 68)
BODY_COLOR = RGBColor(55, 55, 55)
ACCENT_COLOR = RGBColor(0, 112, 192)


def read_csv_rows(path: Path) -> List[Dict[str, str]]:
    with path.open("r", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def read_json(path: Path) -> Dict[str, object]:
    return json.loads(path.read_text(encoding="utf-8"))


def as_float(v: str | None, default: float = 0.0) -> float:
    if v is None:
        return default
    try:
        return float(v)
    except (TypeError, ValueError):
        return default


def find_row(rows: Sequence[Dict[str, str]], **match: str) -> Dict[str, str]:
    for row in rows:
        if all(row.get(k) == v for k, v in match.items()):
            return row
    raise KeyError(f"No row matched {match}")


def set_paragraph_style(paragraph, size: int, bold: bool = False, color: RGBColor | None = None, align=None) -> None:
    paragraph.font.name = FONT_NAME
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    if color is not None:
        paragraph.font.color.rgb = color
    if align is not None:
        paragraph.alignment = align


def set_title(shape, text: str) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    set_paragraph_style(p, 34, bold=True, color=TITLE_COLOR)


def set_lines(tf, lines: Sequence[str], size: int = 20, bullet: bool = True) -> None:
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.text = f"• {line}" if bullet else line
        set_paragraph_style(p, size, color=BODY_COLOR)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    tf_title = slide.shapes.title.text_frame
    tf_title.clear()
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    set_paragraph_style(p_title, 38, bold=True, color=TITLE_COLOR)

    tf_sub = slide.placeholders[1].text_frame
    tf_sub.clear()
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle
    set_paragraph_style(p_sub, 18, color=BODY_COLOR)


def add_bullet_slide(prs: Presentation, title: str, subtitle: str | None, bullets: Sequence[str], size: int = 20) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide.shapes.title, title)
    tf = slide.placeholders[1].text_frame
    tf.clear()
    if subtitle:
        p0 = tf.paragraphs[0]
        p0.text = subtitle
        set_paragraph_style(p0, size, bold=True, color=ACCENT_COLOR)
        p0.space_after = Pt(10)
        for line in bullets:
            p = tf.add_paragraph()
            p.text = f"• {line}"
            set_paragraph_style(p, size, color=BODY_COLOR)
    else:
        set_lines(tf, bullets, size=size, bullet=True)


def add_main_trend_chart(prs: Presentation, metrics: Dict[str, float]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    set_title(slide.shapes.title, "전: Main 성능 추세 (핵심만)")

    chart_data = CategoryChartData()
    chart_data.categories = ["Phase-2", "v5", "v7", "v8 단일", "v8 앙상블"]
    chart_data.add_series(
        "Emotion Macro-F1",
        [
            metrics["phase2_main_f1"],
            metrics["v5_main_f1"],
            metrics["v7_main_f1"],
            metrics["v8_single_f1"],
            metrics["v8_ens_f1"],
        ],
    )

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8),
        Inches(1.5),
        Inches(8.2),
        Inches(4.6),
        chart_data,
    ).chart
    chart.has_legend = False
    chart.value_axis.minimum_scale = 0.0
    chart.value_axis.maximum_scale = 0.8

    msg = slide.shapes.add_textbox(Inches(9.3), Inches(1.7), Inches(3.8), Inches(3.8)).text_frame
    msg.clear()
    lines = [
        f"출발점: {metrics['phase2_main_f1']:.4f}",
        f"현재 최고: {metrics['v8_ens_f1']:.4f}",
        f"단일 최고: {metrics['v8_single_f1']:.4f}",
        "",
        "해석:",
        "표현(HuBERT)+학습레시피+결합 구조가",
        "누적되며 성능이 상승했습니다.",
    ]
    for i, line in enumerate(lines):
        p = msg.paragraphs[0] if i == 0 else msg.add_paragraph()
        p.text = line
        set_paragraph_style(p, 16, bold=(line.startswith("해석")), color=BODY_COLOR, align=PP_ALIGN.LEFT)


def add_cross_chart(prs: Presentation, metrics: Dict[str, float]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    set_title(slide.shapes.title, "전: Cross 일반화 결과 (병목과 개선)")

    chart_data = CategoryChartData()
    chart_data.categories = ["CREMA->RAVDESS", "RAVDESS->CREMA"]
    chart_data.add_series("Phase-2 fusion", [metrics["phase2_c2r_f1"], metrics["phase2_r2c_f1"]])
    chart_data.add_series("v8 + CORAL", [metrics["v8_c2r_f1"], metrics["v8_r2c_f1"]])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8),
        Inches(1.5),
        Inches(8.2),
        Inches(4.6),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.value_axis.minimum_scale = 0.0
    chart.value_axis.maximum_scale = 0.5

    msg = slide.shapes.add_textbox(Inches(9.3), Inches(1.7), Inches(3.8), Inches(3.8)).text_frame
    msg.clear()
    delta_c2r = metrics["v8_c2r_f1"] - metrics["phase2_c2r_f1"]
    delta_r2c = metrics["v8_r2c_f1"] - metrics["phase2_r2c_f1"]
    lines = [
        f"C->R: {metrics['phase2_c2r_f1']:.4f} -> {metrics['v8_c2r_f1']:.4f}",
        f"R->C: {metrics['phase2_r2c_f1']:.4f} -> {metrics['v8_r2c_f1']:.4f}",
        "",
        f"개선폭: +{delta_c2r:.4f} / +{delta_r2c:.4f}",
        "",
        "해석:",
        "도메인 갭은 줄었지만 절대수준은",
        "main 대비 여전히 낮습니다.",
    ]
    for i, line in enumerate(lines):
        p = msg.paragraphs[0] if i == 0 else msg.add_paragraph()
        p.text = line
        set_paragraph_style(p, 16, bold=(line.startswith("해석")), color=BODY_COLOR, align=PP_ALIGN.LEFT)


def load_metrics() -> Dict[str, float]:
    manifest = read_json(MANIFEST_SUMMARY)
    p2 = read_csv_rows(PHASE2_GLOBAL)
    p2_boot = read_csv_rows(PHASE2_BOOTSTRAP)
    p3 = read_csv_rows(PHASE3_GLOBAL)
    v5 = read_csv_rows(PHASE35_V5_MAIN)
    v7 = read_csv_rows(PHASE35_V7)
    v8 = read_csv_rows(PHASE35_V8)
    cross = read_csv_rows(PHASE35_CROSS)

    row_p2_main = find_row(p2, run="main", mode="fusion")
    row_p2_main_audio = find_row(p2, run="main", mode="audio")
    row_p2_main_video = find_row(p2, run="main", mode="video")
    row_p2_c2r = find_row(p2, run="cross_crema_to_ravdess", mode="fusion")
    row_p2_r2c = find_row(p2, run="cross_ravdess_to_crema", mode="fusion")
    row_p2_main_fusion_vs_audio = find_row(p2_boot, run="main", lhs_mode="fusion", rhs_mode="audio")
    row_p2_main_fusion_vs_video = find_row(p2_boot, run="main", lhs_mode="fusion", rhs_mode="video")
    row_p3_main = find_row(p3, run="main")
    row_v5_main = find_row(v5, name="v5_logreg_main")
    row_v7_main = find_row(v7, name="fp32_v7_ce_ls_ws_gated_wide_main")
    row_v8_single = find_row(v8, name="fp32_v8_hubert_gated_wide_tune4")
    row_v8_ens = find_row(v8, name="fp32_v8_hubert_ensemble_vote3_main_t3_t4")
    row_v8_c2r = find_row(cross, name="v8_hubert_logreg_coral_cross_crema_to_ravdess")
    row_v8_r2c = find_row(cross, name="v8_hubert_logreg_coral_cross_ravdess_to_crema")

    totals = manifest.get("totals", {})
    emo6 = manifest.get("emotion6_counts", {})
    return {
        "all_n": float(totals.get("all", 0)),
        "crema_n": float(totals.get("crema_d", 0)),
        "ravdess_n": float(totals.get("ravdess", 0)),
        "av_n": float(totals.get("multimodal_common6_av", 0)),
        "neutral_n": float(emo6.get("neutral", 0)),
        "angry_n": float(emo6.get("angry", 0)),
        "phase2_main_f1": as_float(row_p2_main.get("emotion_macro_f1")),
        "phase2_main_audio_f1": as_float(row_p2_main_audio.get("emotion_macro_f1")),
        "phase2_main_video_f1": as_float(row_p2_main_video.get("emotion_macro_f1")),
        "phase2_main_fusion_vs_audio_delta": as_float(row_p2_main_fusion_vs_audio.get("delta_macro_f1_mean")),
        "phase2_main_fusion_vs_audio_ci_lo": as_float(row_p2_main_fusion_vs_audio.get("delta_macro_f1_ci95_lo")),
        "phase2_main_fusion_vs_audio_ci_hi": as_float(row_p2_main_fusion_vs_audio.get("delta_macro_f1_ci95_hi")),
        "phase2_main_fusion_vs_video_delta": as_float(row_p2_main_fusion_vs_video.get("delta_macro_f1_mean")),
        "phase2_main_fusion_vs_video_ci_lo": as_float(row_p2_main_fusion_vs_video.get("delta_macro_f1_ci95_lo")),
        "phase2_main_fusion_vs_video_ci_hi": as_float(row_p2_main_fusion_vs_video.get("delta_macro_f1_ci95_hi")),
        "phase3_main_f1": as_float(row_p3_main.get("phase3_emotion_macro_f1")),
        "v5_main_f1": as_float(row_v5_main.get("emotion_macro_f1")),
        "v7_main_f1": as_float(row_v7_main.get("emotion_macro_f1")),
        "v8_single_f1": as_float(row_v8_single.get("emotion_macro_f1")),
        "v8_ens_f1": as_float(row_v8_ens.get("emotion_macro_f1")),
        "phase2_c2r_f1": as_float(row_p2_c2r.get("emotion_macro_f1")),
        "phase2_r2c_f1": as_float(row_p2_r2c.get("emotion_macro_f1")),
        "v8_c2r_f1": as_float(row_v8_c2r.get("emotion_macro_f1")),
        "v8_r2c_f1": as_float(row_v8_r2c.get("emotion_macro_f1")),
    }


def build_ppt(metrics: Dict[str, float]) -> None:
    prs = Presentation()

    add_title_slide(
        prs,
        "CREMA-D + RAVDESS 멀티모달 감정인식 연구 진행 보고",
        "스토리 플로우: 문제정의 -> 진단 -> 개입 -> 전환점 -> 결론",
    )

    add_bullet_slide(
        prs,
        "한 장 요약 (스토리 핵심)",
        None,
        [
            "문제정의는 실험 뒤가 아니라 연구 시작 시점에 고정한다",
            "baseline은 문제정의의 근거를 확인하는 진단 단계로 사용한다",
            "개입(표현/학습/도메인적응) 후 main과 cross 모두 개선됐다",
            "결론: 0.7 구간 진입은 확인, 0.9는 별도 전환 트랙이 필요하다",
        ],
        size=22,
    )

    add_bullet_slide(
        prs,
        "발표 흐름 (하나의 이야기)",
        None,
        [
            "1막: 연구 시작 전 문제정의와 가설",
            "2막: baseline 진단과 병목 확인",
            "3막: 병목별 개입(표현, 학습, 도메인적응)",
            "4막: 전환점(HuBERT)과 결과",
            "5막: 결론과 다음 의사결정",
        ],
        size=21,
    )

    add_bullet_slide(
        prs,
        "1막. 연구 시작점: 사전 문제정의와 가설",
        "중요: 문제정의가 먼저이고, baseline은 나중에 검증한다",
        [
            "문제상황: 실제 적용에서는 학습 데이터와 다른 도메인이 반드시 등장한다",
            "목표: in-domain 정확도뿐 아니라 cross-domain 일반화까지 확보",
            "가설 H1: 표현력(pretrained) 강화가 main 성능을 높인다",
            "가설 H2: 도메인적응(CORAL)이 cross 붕괴를 완화한다",
            "가설 H3: 불균형 대응(weighted sampler/label smoothing)이 macro-F1을 안정화한다",
        ],
    )

    add_bullet_slide(
        prs,
        "2막. baseline 진단: 가설을 수치로 점검",
        "여기서부터는 문제정의를 검증하는 단계",
        [
            f"main baseline: audio {metrics['phase2_main_audio_f1']:.4f}, video {metrics['phase2_main_video_f1']:.4f}, fusion {metrics['phase2_main_f1']:.4f}",
            f"fusion-audio Δ={metrics['phase2_main_fusion_vs_audio_delta']:+.4f}, 95%CI[{metrics['phase2_main_fusion_vs_audio_ci_lo']:.4f},{metrics['phase2_main_fusion_vs_audio_ci_hi']:.4f}]",
            f"fusion-video Δ={metrics['phase2_main_fusion_vs_video_delta']:+.4f}, 95%CI[{metrics['phase2_main_fusion_vs_video_ci_lo']:.4f},{metrics['phase2_main_fusion_vs_video_ci_hi']:.4f}]",
            f"cross baseline: C->R {metrics['phase2_c2r_f1']:.4f}, R->C {metrics['phase2_r2c_f1']:.4f}",
            "판단: 멀티모달은 효과가 있으나, cross 붕괴가 핵심 병목",
        ],
    )

    add_bullet_slide(
        prs,
        "3막-A. 개입 1: 입력 표현 강화",
        "무엇/왜/어디/어떻게",
        [
            "무엇: log-mel/delta + pretrained embedding",
            "왜: 저차원 통계특징만으로는 감정 분별 정보가 부족",
            "어디: `scripts/prepare_advanced_features.py`",
            "어떻게: cache_v1 -> v5_hubert 단계 확장, 동일 split에서 비교",
            f"결과 흐름: main F1 {metrics['phase2_main_f1']:.4f} -> {metrics['v5_main_f1']:.4f}",
        ],
    )

    add_bullet_slide(
        prs,
        "3막-B. 개입 2: 학습 안정화",
        "무엇/왜/어디/어떻게",
        [
            "무엇: dual encoder + gated fusion + multitask(emotion/a2/a3)",
            "왜: 모달 상보성 확보 + 불균형/과적합 대응",
            "어디: `scripts/train_fp32_multitask.py`",
            "어떻게: weighted sampler, label smoothing, modality dropout 조합",
            f"결과 흐름: main F1 {metrics['v5_main_f1']:.4f} -> {metrics['v7_main_f1']:.4f}",
        ],
    )

    add_bullet_slide(
        prs,
        "3막-C. 개입 3: 도메인적응",
        "무엇/왜/어디/어떻게",
        [
            "무엇: CORAL(Covariance Alignment)",
            "왜: source/target feature 분포 차이(도메인 갭) 완화",
            "어디: `scripts/train_ml_baselines.py`의 `coral_align_train_to_val`",
            "어떻게: 타깃 라벨 없이 공분산 정렬 후 분류기 학습/평가",
            f"결과: cross C->R {metrics['phase2_c2r_f1']:.4f}->{metrics['v8_c2r_f1']:.4f}, R->C {metrics['phase2_r2c_f1']:.4f}->{metrics['v8_r2c_f1']:.4f}",
        ],
    )

    add_bullet_slide(
        prs,
        "4막. 전환점: HuBERT + FP32 튜닝",
        "개입이 누적되며 성능 점프가 발생한 구간",
        [
            "무엇: audio pretrained를 wav2vec2에서 HuBERT로 전환",
            "왜: 표현 품질을 한 단계 더 올려 단일모델 임계점 접근",
            "어디: `scripts/run_phase35_next_v8_hubert_main.sh`",
            "어떻게: cache_v5_hubert 생성 후 FP32 gated wide 튜닝",
            f"결과: 단일 {metrics['v8_single_f1']:.4f}, 앙상블 {metrics['v8_ens_f1']:.4f}",
        ],
    )

    add_main_trend_chart(prs, metrics)
    add_cross_chart(prs, metrics)

    gap_single = 0.9 - metrics["v8_single_f1"]
    add_bullet_slide(
        prs,
        "5막. 해석: 무엇이 통했고, 무엇이 남았나",
        None,
        [
            "통한 것: 표현강화 + 학습안정화 + 도메인적응의 조합 전략",
            "남은 것: cross 절대성능은 아직 낮고 0.9까지 큰 격차가 존재",
            f"현재 위치: 0.7 구간 진입(앙상블), 단일 0.9까지 +{gap_single:.4f}",
            "핵심 교훈: 문제정의-가설-개입-검증 흐름을 유지해야 다음 실험이 설득된다",
        ],
    )

    add_bullet_slide(
        prs,
        "결론. 다음 의사결정",
        "트랙을 분리해 연구 리스크를 낮춘다",
        [
            f"근거1: cross 최고도 0.3207/0.3187로 main({metrics['v8_ens_f1']:.4f}) 대비 격차 큼",
            "근거2: ROI/품질플래그, partial unfreeze fine-tuning, deep adaptation은 아직 미완료",
            "근거3: 따라서 동일 접근 반복보다 구조 전환형 실험이 필요",
            "트랙 A(단기): 단일모델 0.7 안정 재현(시드/튜닝/분산관리)",
            "트랙 B(중기): 0.9 전환(ROI/품질플래그 + partial unfreeze + deep adaptation)",
            "운영 원칙: 새 실험 시작 전 '문제근거->가설->방법->판정기준' 1페이지 명시",
        ],
    )

    add_bullet_slide(
        prs,
        "마무리",
        None,
        [
            "이 연구는 '많은 실험'보다 '의도와 근거가 연결된 실험'으로 정리했다",
            "핵심 질문에 대한 답: 멀티모달 고도화는 효과가 있었고, 다음 병목도 명확해졌다",
            "질의응답",
        ],
        size=23,
    )

    add_bullet_slide(
        prs,
        "Q&A",
        "답변 원칙: 의도 -> 방법 -> 결과",
        [
            "질문이 오면 먼저 해당 문제정의의 근거 수치를 제시",
            "수치 질문은 발표 슬라이드/요약 파일의 동일 값으로만 답변",
            "설정 질문은 코드 위치와 실험명까지 함께 제시",
            "상세 문답은 확장본 Q&A 문서로 연결",
        ],
        size=20,
    )

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)


def build_pdf(metrics: Dict[str, float]) -> None:
    pdfmetrics.registerFont(TTFont("MalgunGothic", "/mnt/c/Windows/Fonts/malgun.ttf"))
    c = canvas.Canvas(str(OUT_PDF), pagesize=A4)
    w, h = A4
    y = h - 20 * mm

    def line(text: str, size: int = 11, leading: float = 6.0, bold: bool = False) -> None:
        nonlocal y
        font = "MalgunGothic"
        c.setFont(font, size)
        c.drawString(18 * mm, y, text)
        y -= (size + leading)

    line("연구 진행 요약(스토리텔링 보강판)", size=16, bold=True)
    line("기준: 2026-02-24 결과", size=10)
    y -= 4

    line("[1막] 문제정의와 가설(먼저 정의)", size=13)
    line("- 문제정의는 baseline 이후가 아니라 연구 시작 시점에 고정", size=11)
    line("- 목표: main 성능 + cross 일반화 동시 개선", size=11)
    y -= 2

    line("[2막] baseline 진단(가설 검증 시작)", size=13)
    line(
        f"- main audio/video/fusion: {metrics['phase2_main_audio_f1']:.4f}/{metrics['phase2_main_video_f1']:.4f}/{metrics['phase2_main_f1']:.4f}",
        size=11,
    )
    line(f"- cross: {metrics['phase2_c2r_f1']:.4f}/{metrics['phase2_r2c_f1']:.4f}", size=11)
    y -= 2

    line("[3막] 개입(무엇/왜/어디/어떻게)", size=13)
    line("- 표현강화(cache 고도화, HuBERT), 학습안정화(gated/샘플링/스무딩), CORAL", size=11)
    line("- 코드 기반 추적: prepare_advanced_features.py / train_fp32_multitask.py / train_ml_baselines.py", size=11)
    y -= 2

    line("[4막] 결과", size=13)
    line(
        f"- main F1: {metrics['phase2_main_f1']:.4f} -> {metrics['v8_ens_f1']:.4f} (앙상블), 단일 {metrics['v8_single_f1']:.4f}",
        size=11,
    )
    line(
        f"- cross F1: C->R {metrics['phase2_c2r_f1']:.4f}->{metrics['v8_c2r_f1']:.4f}, R->C {metrics['phase2_r2c_f1']:.4f}->{metrics['v8_r2c_f1']:.4f}",
        size=11,
    )
    y -= 2

    line("[5막] 결론과 다음 단계", size=13)
    line("- 0.7 구간 진입은 확인, 0.9는 구조 전환형 고도화 필요", size=11)
    line("- 다음 단계: (A) 0.7 안정화 (B) 0.9 전환 트랙 분리", size=11)

    c.showPage()
    c.save()


def main() -> None:
    metrics = load_metrics()
    build_ppt(metrics)
    build_pdf(metrics)
    print(f"[OK] optimized ppt: {OUT_PPTX}")
    print(f"[OK] optimized pdf: {OUT_PDF}")


if __name__ == "__main__":
    main()
