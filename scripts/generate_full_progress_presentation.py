#!/usr/bin/env python3
"""한국어 연구 진행 발표자료(PPTX/PDF) 생성기 (상세 보강판)."""

from __future__ import annotations

import csv
import json
from pathlib import Path
from typing import Dict, List, Sequence, Tuple

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas


ROOT = Path(".")
REPORTS = ROOT / "derived" / "reports"
MANIFEST_SUMMARY = ROOT / "derived" / "manifests" / "summary.json"
RUNTIME_BENCH = ROOT / "derived" / "results" / "phase2_runtime_bench_fusion_fold0.json"
CACHE_V2_SUMMARY = ROOT / "derived" / "features" / "cache_v2" / "summary.json"
CACHE_V3_SUMMARY = ROOT / "derived" / "features" / "cache_v3" / "summary.json"
CACHE_V4_SUMMARY = ROOT / "derived" / "features" / "cache_v4" / "summary.json"
CACHE_V5_SUMMARY = ROOT / "derived" / "features" / "cache_v5_hubert" / "summary.json"

PHASE2_GLOBAL = REPORTS / "phase2_global_metrics.csv"
PHASE3_GLOBAL = REPORTS / "phase3_global_metrics.csv"
PHASE2_BOOTSTRAP = REPORTS / "phase2_pairwise_bootstrap.csv"
PHASE3_BOOTSTRAP = REPORTS / "phase3_vs_phase2_bootstrap.csv"
PHASE35_V2 = REPORTS / "phase35_advancement_metrics.csv"
PHASE35_V3 = REPORTS / "phase35_advancement_v2_main_metrics.csv"
PHASE35_V5 = REPORTS / "phase35_strong_v1_main_metrics.csv"
PHASE35_V6 = REPORTS / "phase35_next_v6_metrics.csv"
PHASE35_V7 = REPORTS / "phase35_next_v7_metrics.csv"
PHASE35_V8 = REPORTS / "phase35_next_v8_metrics.csv"
PHASE35_CROSS = REPORTS / "phase35_cross_domain_adapt_metrics.csv"
V7_WIDE_SUMMARY = ROOT / "derived" / "results" / "fp32_multitask_phase35_v7_ce_ls_ws_gated_wide_main" / "summary.json"
V8_TUNE4_SUMMARY = ROOT / "derived" / "results" / "fp32_multitask_phase35_v8_hubert_gated_wide_tune4" / "summary.json"

CROSS_SPLIT_FILES = [
    "test_crema_train_ravdess_common6_av_train.txt",
    "test_crema_train_ravdess_common6_av_test.txt",
    "train_crema_test_ravdess_common6_av_train.txt",
    "train_crema_test_ravdess_common6_av_test.txt",
    "test_crema_train_ravdess_common6_train.txt",
    "test_crema_train_ravdess_common6_test.txt",
    "train_crema_test_ravdess_common6_train.txt",
    "train_crema_test_ravdess_common6_test.txt",
]

OUT_PPTX = REPORTS / "research_progress_summary_2026-02-24.pptx"
OUT_PDF = REPORTS / "research_progress_summary_2026-02-24.pdf"

KOR_FONT = "Malgun Gothic"
TITLE_COLOR = RGBColor(14, 32, 63)
SUB_COLOR = RGBColor(66, 66, 66)
ACCENT_COLOR = RGBColor(0, 114, 198)
SECTION_BG = RGBColor(18, 47, 91)
CARD_BG_A = RGBColor(239, 246, 255)
CARD_BG_B = RGBColor(243, 250, 244)
CARD_BG_C = RGBColor(255, 245, 235)
CARD_BG_D = RGBColor(247, 242, 255)


# ========================
# Data helpers
# ========================


def read_csv_rows(path: Path) -> List[Dict[str, str]]:
    with path.open("r", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def read_json(path: Path, default=None):
    if not path.exists():
        return {} if default is None else default
    return json.loads(path.read_text(encoding="utf-8"))


def as_float(v: str | None, default: float = 0.0) -> float:
    if v is None:
        return default
    try:
        return float(v)
    except (TypeError, ValueError):
        return default


def count_lines(path: Path) -> int:
    if not path.exists():
        return 0
    with path.open("r", encoding="utf-8") as f:
        return sum(1 for _ in f)


def find_row(rows: Sequence[Dict[str, str]], key: str, value: str) -> Dict[str, str]:
    for r in rows:
        if r.get(key) == value:
            return r
    raise KeyError(f"Missing row where {key}={value!r}")


def find_phase2(rows: Sequence[Dict[str, str]], run: str, mode: str) -> Dict[str, str]:
    for r in rows:
        if r.get("run") == run and r.get("mode") == mode:
            return r
    raise KeyError(f"Missing phase2 row run={run}, mode={mode}")


def row_get(rows: Sequence[Dict[str, str]], key: str, value: str) -> Dict[str, str]:
    for r in rows:
        if r.get(key) == value:
            return r
    return {}


def format_ci(lo: float, hi: float) -> str:
    return f"[{lo:.4f}, {hi:.4f}]"


# ========================
# PPT helpers
# ========================


def style_paragraph(paragraph, size: int = 18, bold: bool = False, color: RGBColor | None = None, align=None) -> None:
    paragraph.font.name = KOR_FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    if color is not None:
        paragraph.font.color.rgb = color
    if align is not None:
        paragraph.alignment = align


def fill_shape(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_title(shape, text: str, size: int = 30) -> None:
    shape.text = text
    style_paragraph(shape.text_frame.paragraphs[0], size=size, bold=True, color=TITLE_COLOR)


def set_text_lines(text_frame, lines: List[str], size: int = 18, bullet: bool = True) -> None:
    text_frame.clear()
    for idx, line in enumerate(lines):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.level = 0
        p.text = f"• {line}" if bullet else line
        style_paragraph(p, size=size, color=SUB_COLOR)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    t = slide.shapes.title.text_frame
    t.text = title
    style_paragraph(t.paragraphs[0], size=38, bold=True, color=TITLE_COLOR)

    s = slide.placeholders[1].text_frame
    s.text = subtitle
    style_paragraph(s.paragraphs[0], size=18, color=SUB_COLOR)


def add_section_divider(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill_shape(bg, SECTION_BG)
    bg.line.fill.background()

    t_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.8), Inches(1.2))
    t_tf = t_box.text_frame
    t_tf.text = title
    style_paragraph(t_tf.paragraphs[0], size=42, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.LEFT)

    s_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.8), Inches(1.4))
    s_tf = s_box.text_frame
    s_tf.text = subtitle
    style_paragraph(s_tf.paragraphs[0], size=20, color=RGBColor(225, 235, 250), align=PP_ALIGN.LEFT)


def add_agenda_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide.shapes.title, "발표 목차")

    lines = [
        "I. 연구 맥락: 문제정의, KPI, 데이터/분할 프로토콜",
        "II. 방법론: 전처리/학습 이론, 세부 설정값, 코드 적용 위치",
        "III. 단계별 실행: Phase-1/2/3/3.5의 가설-방법-결과-인사이트",
        "IV. 종합 분석: 성능 추세, ablation, 설계 대비 구현 현황",
        "V. 의사결정: 0.7 달성 해석과 0.9 전환 로드맵",
    ]
    set_text_lines(slide.shapes.placeholders[1].text_frame, lines, size=20, bullet=False)


def add_bullet_slide(prs: Presentation, title: str, subtitle: str | None, bullets: List[str], size: int = 18) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide.shapes.title, title)

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    if subtitle:
        p0 = tf.paragraphs[0]
        p0.text = subtitle
        style_paragraph(p0, size=size + 1, bold=True, color=ACCENT_COLOR)
        p0.space_after = Pt(10)

    for i, line in enumerate(bullets):
        p = tf.add_paragraph() if (subtitle or i > 0) else tf.paragraphs[0]
        p.text = f"• {line}"
        style_paragraph(p, size=size, color=SUB_COLOR)


def add_kpi_cards_slide(prs: Presentation, title: str, cards: List[Tuple[str, str, str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(slide.shapes.title, title)

    lefts = [0.45, 4.55, 8.65]
    colors = [CARD_BG_A, CARD_BG_B, CARD_BG_C]

    for idx, (header, value, desc) in enumerate(cards):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(lefts[idx]),
            Inches(1.6),
            Inches(4.0),
            Inches(5.0),
        )
        fill_shape(card, colors[idx])
        card.line.color.rgb = RGBColor(210, 220, 235)

        tf = card.text_frame
        tf.clear()

        p0 = tf.paragraphs[0]
        p0.text = header
        style_paragraph(p0, size=18, bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)

        p1 = tf.add_paragraph()
        p1.text = value
        style_paragraph(p1, size=30, bold=True, color=ACCENT_COLOR, align=PP_ALIGN.CENTER)

        p2 = tf.add_paragraph()
        p2.text = desc
        style_paragraph(p2, size=14, color=SUB_COLOR, align=PP_ALIGN.CENTER)


def add_table_slide(
    prs: Presentation,
    title: str,
    columns: List[str],
    rows: List[List[str]],
    col_widths: List[float] | None = None,
    top: float = 1.2,
    height: float = 5.7,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(slide.shapes.title, title)

    table = slide.shapes.add_table(
        len(rows) + 1,
        len(columns),
        Inches(0.35),
        Inches(top),
        Inches(12.6),
        Inches(height),
    ).table

    if col_widths and len(col_widths) == len(columns):
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(width)

    for i, col in enumerate(columns):
        cell = table.cell(0, i)
        cell.text = col
        style_paragraph(cell.text_frame.paragraphs[0], size=13, bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            align = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
            style_paragraph(cell.text_frame.paragraphs[0], size=12, color=SUB_COLOR, align=align)


def add_2x2_blocks_slide(prs: Presentation, title: str, blocks: List[Tuple[str, List[str]]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    t_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(12.5), Inches(0.7))
    t_tf = t_box.text_frame
    t_tf.text = title
    style_paragraph(t_tf.paragraphs[0], size=30, bold=True, color=TITLE_COLOR)

    pos = [
        (0.45, 1.1, 6.1, 2.8),
        (6.75, 1.1, 6.1, 2.8),
        (0.45, 4.05, 6.1, 2.8),
        (6.75, 4.05, 6.1, 2.8),
    ]
    fills = [CARD_BG_A, CARD_BG_B, CARD_BG_C, CARD_BG_D]

    for idx, (heading, lines) in enumerate(blocks[:4]):
        l, t, w, h = pos[idx]
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        fill_shape(box, fills[idx])
        box.line.color.rgb = RGBColor(206, 216, 232)

        tf = box.text_frame
        tf.clear()

        p0 = tf.paragraphs[0]
        p0.text = heading
        style_paragraph(p0, size=18, bold=True, color=TITLE_COLOR)

        for line in lines:
            p = tf.add_paragraph()
            p.text = f"• {line}"
            style_paragraph(p, size=14, color=SUB_COLOR)


def add_main_trend_chart(prs: Presentation, m: Dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(slide.shapes.title, "Main 성능 추이 (Emotion Macro-F1)")

    categories = ["P2", "P3", "v2", "v3", "v5", "v6", "v7", "v8 단일", "v8 앙상블"]
    values = [
        m["phase2_main_f1"],
        m["phase3_main_f1"],
        m["v2_main_f1"],
        m["v3_main_f1"],
        m["v5_main_f1"],
        m["v6_main_f1"],
        m["v7_main_f1"],
        m["v8_single_f1"],
        m["v8_ensemble_f1"],
    ]

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Macro-F1", values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(0.6),
        Inches(1.4),
        Inches(12.1),
        Inches(4.8),
        chart_data,
    ).chart

    chart.has_legend = False
    chart.value_axis.minimum_scale = 0.0
    chart.value_axis.maximum_scale = 0.8
    chart.value_axis.has_major_gridlines = True

    msg = (
        f"Phase-2(0.3950) -> v8 단일({m['v8_single_f1']:.4f}) -> v8 앙상블({m['v8_ensemble_f1']:.4f}) | "
        "표현+결합+적응의 누적 효과"
    )
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.25), Inches(12.0), Inches(0.7))
    tf = box.text_frame
    tf.text = msg
    style_paragraph(tf.paragraphs[0], size=15, bold=True, color=ACCENT_COLOR)


def add_cross_chart(prs: Presentation, m: Dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(slide.shapes.title, "Cross-domain 성능 비교 (Macro-F1)")

    chart_data = CategoryChartData()
    chart_data.categories = ["CREMA->RAVDESS", "RAVDESS->CREMA"]
    chart_data.add_series("Phase-2 fusion", [m["phase2_cross_c2r_f1"], m["phase2_cross_r2c_f1"]])
    chart_data.add_series("v5 logreg+CORAL", [m["v5_cross_c2r_f1"], m["v5_cross_r2c_f1"]])
    chart_data.add_series("v8 hubert+CORAL", [m["v8_cross_c2r_f1"], m["v8_cross_r2c_f1"]])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.6),
        Inches(1.4),
        Inches(12.1),
        Inches(4.8),
        chart_data,
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.minimum_scale = 0.0
    chart.value_axis.maximum_scale = 0.4

    note = (
        f"v8 최종 cross: C->R {m['v8_cross_c2r_f1']:.4f}, R->C {m['v8_cross_r2c_f1']:.4f} | "
        "CORAL + HuBERT 결합이 양방향 개선"
    )
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.25), Inches(12.0), Inches(0.7))
    tf = box.text_frame
    tf.text = note
    style_paragraph(tf.paragraphs[0], size=15, bold=True, color=ACCENT_COLOR)


def add_phase2_modality_chart(prs: Presentation, m: Dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(slide.shapes.title, "Phase-2 모달리티별 성능 비교 (Macro-F1)")

    p2_main = {r["mode"]: r for r in m["phase2_main_modes"]}
    p2_c2r = {r["mode"]: r for r in m["phase2_cross_modes"]["cross_crema_to_ravdess"]}
    p2_r2c = {r["mode"]: r for r in m["phase2_cross_modes"]["cross_ravdess_to_crema"]}

    chart_data = CategoryChartData()
    chart_data.categories = ["main", "C->R", "R->C"]
    chart_data.add_series("audio", [p2_main["audio"]["f1"], p2_c2r["audio"]["f1"], p2_r2c["audio"]["f1"]])
    chart_data.add_series("video", [p2_main["video"]["f1"], p2_c2r["video"]["f1"], p2_r2c["video"]["f1"]])
    chart_data.add_series("fusion", [p2_main["fusion"]["f1"], p2_c2r["fusion"]["f1"], p2_r2c["fusion"]["f1"]])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.6),
        Inches(1.4),
        Inches(12.1),
        Inches(4.8),
        chart_data,
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.minimum_scale = 0.0
    chart.value_axis.maximum_scale = 0.45

    note = "해석: main에서는 fusion 우세, cross에서는 방향별로 audio/fusion 우세가 갈리는 도메인 비대칭 확인"
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.25), Inches(12.0), Inches(0.7))
    tf = box.text_frame
    tf.text = note
    style_paragraph(tf.paragraphs[0], size=15, bold=True, color=ACCENT_COLOR)


def add_phase_card(
    prs: Presentation,
    title: str,
    question: str,
    methods: List[str],
    comparison: List[str],
    insight: List[str],
) -> None:
    add_2x2_blocks_slide(
        prs,
        title,
        [
            ("연구질문", [question]),
            ("방법론", methods),
            ("비교/평가", comparison),
            ("결론/인사이트", insight),
        ],
    )


# ========================
# Metric loading
# ========================


def load_metrics() -> Dict[str, object]:
    phase2 = read_csv_rows(PHASE2_GLOBAL)
    phase3 = read_csv_rows(PHASE3_GLOBAL)
    phase2_boot = read_csv_rows(PHASE2_BOOTSTRAP)
    phase3_boot = read_csv_rows(PHASE3_BOOTSTRAP)
    v2 = read_csv_rows(PHASE35_V2)
    v3 = read_csv_rows(PHASE35_V3)
    v5 = read_csv_rows(PHASE35_V5)
    v6 = read_csv_rows(PHASE35_V6)
    v7 = read_csv_rows(PHASE35_V7)
    v8 = read_csv_rows(PHASE35_V8)
    cross = read_csv_rows(PHASE35_CROSS)

    manifest = read_json(MANIFEST_SUMMARY, default={})
    runtime = read_json(RUNTIME_BENCH, default={})
    cache_v2 = read_json(CACHE_V2_SUMMARY, default={})
    cache_v3 = read_json(CACHE_V3_SUMMARY, default={})
    cache_v4 = read_json(CACHE_V4_SUMMARY, default={})
    cache_v5 = read_json(CACHE_V5_SUMMARY, default={})
    v7_wide_summary = read_json(V7_WIDE_SUMMARY, default={})
    v8_tune4_summary = read_json(V8_TUNE4_SUMMARY, default={})

    out: Dict[str, object] = {}

    p2_main_fusion = find_phase2(phase2, "main", "fusion")
    p2_main_audio = find_phase2(phase2, "main", "audio")
    p2_main_video = find_phase2(phase2, "main", "video")

    p3_main = find_row(phase3, "run", "main")
    p3_c2r = find_row(phase3, "run", "cross_crema_to_ravdess")
    p3_r2c = find_row(phase3, "run", "cross_ravdess_to_crema")

    v2_best = find_row(v2, "name", "ml_v2_logreg_fusion")
    v3_best = find_row(v3, "name", "ml_v3_rbfsvm_fusion")
    v5_best = find_row(v5, "name", "v5_logreg_main")
    v6_best = find_row(v6, "name", "fp32_v6_ce_ls_ws_main")
    v7_best = find_row(v7, "name", "fp32_v7_ce_ls_ws_gated_wide_main")
    v8_single = find_row(v8, "name", "fp32_v8_hubert_gated_wide_tune4")
    v8_ens = find_row(v8, "name", "fp32_v8_hubert_ensemble_vote3_main_t3_t4")

    out["phase2_main_f1"] = as_float(p2_main_fusion.get("emotion_macro_f1"))
    out["phase3_main_f1"] = as_float(p3_main.get("phase3_emotion_macro_f1"))
    out["v2_main_f1"] = as_float(v2_best.get("emotion_macro_f1"))
    out["v3_main_f1"] = as_float(v3_best.get("emotion_macro_f1"))
    out["v5_main_f1"] = as_float(v5_best.get("emotion_macro_f1"))
    out["v6_main_f1"] = as_float(v6_best.get("emotion_macro_f1"))
    out["v7_main_f1"] = as_float(v7_best.get("emotion_macro_f1"))
    out["v8_single_f1"] = as_float(v8_single.get("emotion_macro_f1"))
    out["v8_ensemble_f1"] = as_float(v8_ens.get("emotion_macro_f1"))

    out["phase2_cross_c2r_f1"] = as_float(find_phase2(phase2, "cross_crema_to_ravdess", "fusion").get("emotion_macro_f1"))
    out["phase2_cross_r2c_f1"] = as_float(find_phase2(phase2, "cross_ravdess_to_crema", "fusion").get("emotion_macro_f1"))
    out["v5_cross_c2r_f1"] = as_float(find_row(cross, "name", "v5_logreg_coral_cross_crema_to_ravdess").get("emotion_macro_f1"))
    out["v5_cross_r2c_f1"] = as_float(find_row(cross, "name", "v5_logreg_coral_cross_ravdess_to_crema").get("emotion_macro_f1"))
    out["v8_cross_c2r_f1"] = as_float(find_row(cross, "name", "v8_hubert_logreg_coral_cross_crema_to_ravdess").get("emotion_macro_f1"))
    out["v8_cross_r2c_f1"] = as_float(find_row(cross, "name", "v8_hubert_logreg_coral_cross_ravdess_to_crema").get("emotion_macro_f1"))

    out["phase2_main_modes"] = [
        {
            "mode": "audio",
            "f1": as_float(p2_main_audio.get("emotion_macro_f1")),
            "acc": as_float(p2_main_audio.get("emotion_acc")),
            "ar2": as_float(p2_main_audio.get("arousal2_mae")),
            "ar3": as_float(p2_main_audio.get("arousal3_mae")),
        },
        {
            "mode": "video",
            "f1": as_float(p2_main_video.get("emotion_macro_f1")),
            "acc": as_float(p2_main_video.get("emotion_acc")),
            "ar2": as_float(p2_main_video.get("arousal2_mae")),
            "ar3": as_float(p2_main_video.get("arousal3_mae")),
        },
        {
            "mode": "fusion",
            "f1": as_float(p2_main_fusion.get("emotion_macro_f1")),
            "acc": as_float(p2_main_fusion.get("emotion_acc")),
            "ar2": as_float(p2_main_fusion.get("arousal2_mae")),
            "ar3": as_float(p2_main_fusion.get("arousal3_mae")),
        },
    ]

    out["phase2_cross_modes"] = {
        "cross_crema_to_ravdess": [
            {
                "mode": mode,
                "f1": as_float(find_phase2(phase2, "cross_crema_to_ravdess", mode).get("emotion_macro_f1")),
                "acc": as_float(find_phase2(phase2, "cross_crema_to_ravdess", mode).get("emotion_acc")),
            }
            for mode in ["audio", "video", "fusion"]
        ],
        "cross_ravdess_to_crema": [
            {
                "mode": mode,
                "f1": as_float(find_phase2(phase2, "cross_ravdess_to_crema", mode).get("emotion_macro_f1")),
                "acc": as_float(find_phase2(phase2, "cross_ravdess_to_crema", mode).get("emotion_acc")),
            }
            for mode in ["audio", "video", "fusion"]
        ],
    }

    out["phase3_cross"] = {
        "cross_crema_to_ravdess": as_float(p3_c2r.get("phase3_emotion_macro_f1")),
        "cross_ravdess_to_crema": as_float(p3_r2c.get("phase3_emotion_macro_f1")),
    }

    # Bootstrap summaries
    phase3_boot_map = {r.get("run", ""): r for r in phase3_boot}
    out["phase3_bootstrap"] = {
        run: {
            "delta": as_float(row.get("delta_mean_phase3_minus_phase2fusion")),
            "ci_lo": as_float(row.get("ci95_low")),
            "ci_hi": as_float(row.get("ci95_high")),
            "n": int(as_float(row.get("n_aligned"))),
        }
        for run, row in phase3_boot_map.items()
    }

    out["phase2_pairwise_main"] = [
        {
            "lhs": r.get("lhs_mode", ""),
            "rhs": r.get("rhs_mode", ""),
            "delta": as_float(r.get("delta_macro_f1_mean")),
            "ci_lo": as_float(r.get("delta_macro_f1_ci95_lo")),
            "ci_hi": as_float(r.get("delta_macro_f1_ci95_hi")),
        }
        for r in phase2_boot
        if r.get("run") == "main"
    ]

    # Cross ablation table
    out["cross_ablation"] = {
        "crema_to_ravdess": {
            "phase2_fusion": as_float(find_phase2(phase2, "cross_crema_to_ravdess", "fusion").get("emotion_macro_f1")),
            "v3_baseline": as_float(find_row(cross, "name", "v3_logreg_baseline_cross_crema_to_ravdess").get("emotion_macro_f1")),
            "v5_baseline": as_float(find_row(cross, "name", "v5_logreg_baseline_cross_crema_to_ravdess").get("emotion_macro_f1")),
            "v5_coral": as_float(find_row(cross, "name", "v5_logreg_coral_cross_crema_to_ravdess").get("emotion_macro_f1")),
            "v8_coral": as_float(find_row(cross, "name", "v8_hubert_logreg_coral_cross_crema_to_ravdess").get("emotion_macro_f1")),
        },
        "ravdess_to_crema": {
            "phase2_fusion": as_float(find_phase2(phase2, "cross_ravdess_to_crema", "fusion").get("emotion_macro_f1")),
            "v3_baseline": as_float(find_row(cross, "name", "v3_logreg_baseline_cross_ravdess_to_crema").get("emotion_macro_f1")),
            "v5_baseline": as_float(find_row(cross, "name", "v5_logreg_baseline_cross_ravdess_to_crema").get("emotion_macro_f1")),
            "v5_coral": as_float(find_row(cross, "name", "v5_logreg_coral_cross_ravdess_to_crema").get("emotion_macro_f1")),
            "v8_coral": as_float(find_row(cross, "name", "v8_hubert_logreg_coral_cross_ravdess_to_crema").get("emotion_macro_f1")),
        },
    }

    # Main representative trend with MAE
    out["main_stage_metrics"] = [
        {
            "stage": "Phase-2 fusion",
            "f1": as_float(p2_main_fusion.get("emotion_macro_f1")),
            "acc": as_float(p2_main_fusion.get("emotion_acc")),
            "ar2": as_float(p2_main_fusion.get("arousal2_mae")),
            "ar3": as_float(p2_main_fusion.get("arousal3_mae")),
        },
        {
            "stage": "Phase-3 FP32",
            "f1": as_float(p3_main.get("phase3_emotion_macro_f1")),
            "acc": as_float(p3_main.get("phase3_emotion_acc")),
            "ar2": as_float(p3_main.get("phase3_arousal2_mae")),
            "ar3": as_float(p3_main.get("phase3_arousal3_mae")),
        },
        {
            "stage": "v2 logreg",
            "f1": as_float(v2_best.get("emotion_macro_f1")),
            "acc": as_float(v2_best.get("emotion_acc")),
            "ar2": as_float(v2_best.get("arousal2_mae")),
            "ar3": as_float(v2_best.get("arousal3_mae")),
        },
        {
            "stage": "v3 rbfsvm",
            "f1": as_float(v3_best.get("emotion_macro_f1")),
            "acc": as_float(v3_best.get("emotion_acc")),
            "ar2": as_float(v3_best.get("arousal2_mae")),
            "ar3": as_float(v3_best.get("arousal3_mae")),
        },
        {
            "stage": "v5 logreg",
            "f1": as_float(v5_best.get("emotion_macro_f1")),
            "acc": as_float(v5_best.get("emotion_acc")),
            "ar2": as_float(v5_best.get("arousal2_mae")),
            "ar3": as_float(v5_best.get("arousal3_mae")),
        },
        {
            "stage": "v6 fp32",
            "f1": as_float(v6_best.get("emotion_macro_f1")),
            "acc": as_float(v6_best.get("emotion_acc")),
            "ar2": as_float(v6_best.get("arousal2_mae")),
            "ar3": as_float(v6_best.get("arousal3_mae")),
        },
        {
            "stage": "v7 fp32",
            "f1": as_float(v7_best.get("emotion_macro_f1")),
            "acc": as_float(v7_best.get("emotion_acc")),
            "ar2": as_float(v7_best.get("arousal2_mae")),
            "ar3": as_float(v7_best.get("arousal3_mae")),
        },
        {
            "stage": "v8 단일",
            "f1": as_float(v8_single.get("emotion_macro_f1")),
            "acc": as_float(v8_single.get("emotion_acc")),
            "ar2": as_float(v8_single.get("arousal2_mae")),
            "ar3": as_float(v8_single.get("arousal3_mae")),
        },
        {
            "stage": "v8 앙상블",
            "f1": as_float(v8_ens.get("emotion_macro_f1")),
            "acc": as_float(v8_ens.get("emotion_acc")),
            "ar2": as_float(v8_ens.get("arousal2_mae")),
            "ar3": as_float(v8_ens.get("arousal3_mae")),
        },
    ]

    # Split file counts
    split_dir = ROOT / "derived" / "splits" / "cross_dataset"
    out["cross_split_counts"] = [
        {
            "file": name,
            "count": count_lines(split_dir / name),
            "path": f"derived/splits/cross_dataset/{name}",
        }
        for name in CROSS_SPLIT_FILES
    ]

    top_v8 = sorted(v8, key=lambda r: as_float(r.get("emotion_macro_f1")), reverse=True)[:8]
    out["top_v8"] = top_v8

    out["totals"] = manifest.get("totals", {})
    out["emotion6_counts"] = manifest.get("emotion6_counts", {})
    out["splits_summary"] = manifest.get("splits", {})
    out["runtime"] = runtime
    out["cache_v2"] = cache_v2
    out["cache_v3"] = cache_v3
    out["cache_v4"] = cache_v4
    out["cache_v5"] = cache_v5
    out["v7_wide_run"] = v7_wide_summary.get("run", {})
    out["v8_tune4_run"] = v8_tune4_summary.get("run", {})

    return out


# ========================
# PPT build
# ========================


def build_pptx(m: Dict[str, object]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    totals = m["totals"]
    emotion_counts = m["emotion6_counts"]
    split_summary = m["splits_summary"]
    runtime = m["runtime"]
    cache_v2 = m["cache_v2"]
    cache_v3 = m["cache_v3"]
    cache_v4 = m["cache_v4"]
    cache_v5 = m["cache_v5"]
    v7_wide_run = m["v7_wide_run"]
    v8_tune4_run = m["v8_tune4_run"]

    add_title_slide(
        prs,
        "CREMA-D + RAVDESS 멀티모달 연구 진행 보고",
        "설계서·진행문서·실험 산출물 교차검증 기반 상세 발표자료 (기준일: 2026-02-24)",
    )
    add_agenda_slide(prs)

    # Section I
    add_section_divider(prs, "I. 연구 맥락", "문제정의, 연구질문, KPI, 데이터/평가 프로토콜")

    add_bullet_slide(
        prs,
        "연구 배경과 문제 정의",
        "멀티모달 관점에서 해결해야 했던 핵심 문제",
        [
            "감정 단서는 음성(억양/에너지/리듬)과 영상(표정/동작/시선)에 분산되어 단일모달이 취약합니다.",
            "CREMA-D/RAVDESS 간 도메인 차이로 cross-domain 성능 붕괴가 반복적으로 관찰되었습니다.",
            "그래서 본 연구는 main 향상과 cross 회복을 동시에 달성하는 이중 KPI 구조로 운영되었습니다.",
            "연구 설계는 deep-research-report(기본) + advancement(고도화) 문서 체계를 기준으로 확장되었습니다.",
        ],
    )

    gap_single = 0.9 - float(m["v8_single_f1"])
    gap_ens = 0.92 - float(m["v8_ensemble_f1"])
    add_kpi_cards_slide(
        prs,
        "핵심 KPI 스냅샷",
        [
            ("0.7 트랙(현재)", f"{m['v8_ensemble_f1']:.4f}", "앙상블 기준 달성, 단일모델 0.6992"),
            ("0.9 트랙 격차", f"{gap_single:+.4f}", "단일목표 0.9 대비 부족분"),
            ("Cross 양방향 최고", f"{m['v8_cross_c2r_f1']:.4f} / {m['v8_cross_r2c_f1']:.4f}", "v8 hubert + CORAL"),
        ],
    )

    add_table_slide(
        prs,
        "연구 데이터셋 규모와 모집단",
        ["항목", "수치", "해석"],
        [
            ["전체 샘플", f"{totals.get('all', 0):,}", "CREMA-D + RAVDESS 전체"],
            ["CREMA-D", f"{totals.get('crema_d', 0):,}", "배우 91, 감정 발화 중심"],
            ["RAVDESS", f"{totals.get('ravdess', 0):,}", "배우 24, audio/video 포함"],
            ["Common-6", f"{totals.get('common6_all', 0):,}", "공통 라벨 정합 후 실험 가능 데이터"],
            ["멀티모달 AV(Common-6)", f"{totals.get('multimodal_common6_av', 0):,}", "주 학습/평가셋"],
            ["RAVDESS AO(Common-6)", f"{totals.get('ravdess_audio_only_common6', 0):,}", "오디오 보조 비교셋"],
        ],
        col_widths=[3.7, 2.0, 6.5],
    )

    add_table_slide(
        prs,
        "라벨 분포(Emotion Common-6)",
        ["라벨", "샘플 수", "비율(%)"],
        [
            [label, f"{count:,}", f"{(count / max(1, totals.get('common6_all', 1))) * 100:.2f}"]
            for label, count in sorted(emotion_counts.items(), key=lambda x: x[0])
        ],
        col_widths=[3.6, 3.0, 3.0],
        top=1.4,
        height=4.8,
    )

    split_rows = []
    for key in ["all", "crema_d", "ravdess"]:
        s = split_summary.get(key, {})
        fold_counts = s.get("fold_counts", {})
        fold_text = ", ".join([f"{k}:{v}" for k, v in sorted(fold_counts.items(), key=lambda x: x[0])])
        split_rows.append([key, str(s.get("num_rows", 0)), str(s.get("num_groups", 0)), fold_text])
    add_table_slide(
        prs,
        "main 평가 분할 프로토콜 (GroupKFold)",
        ["split", "rows", "groups(actor)", "fold counts"],
        split_rows,
        col_widths=[1.8, 1.6, 2.0, 7.0],
    )

    cross_rows = [
        [row["path"], str(row["count"]), "cross split 파일" if "common6_" in row["file"] else "" ]
        for row in m["cross_split_counts"]
    ]
    add_table_slide(
        prs,
        "cross 평가 분할 파일(실제 산출물)",
        ["파일 경로", "라인 수", "비고"],
        cross_rows,
        col_widths=[8.6, 1.6, 2.1],
        top=1.3,
        height=5.4,
    )

    cache_rows = [
        ["cache dir", str(cache_v5.get("cache_dir", "-"))],
        ["audio backbone", str(cache_v5.get("audio_pretrained_backbone", "-"))],
        ["device", f"{cache_v5.get('device_requested', '-')}/{cache_v5.get('device_resolved', '-')}"],
        ["cuda", str(cache_v5.get("cuda_name", "-"))],
        ["rows(ok/fail)", f"{cache_v5.get('ok', 0)} / {cache_v5.get('fail', 0)}"],
        ["elapsed_sec", f"{as_float(cache_v5.get('elapsed_sec')):.1f}"],
    ]
    add_table_slide(
        prs,
        "pretrained feature cache 실제 생성 결과 (cache_v5_hubert)",
        ["항목", "값"],
        cache_rows,
        col_widths=[4.0, 8.0],
        top=1.5,
        height=4.8,
    )

    # Section II
    add_section_divider(prs, "II. 방법론", "전처리·모델·평가를 단계별로 통제해 성능 기여를 분해")

    add_2x2_blocks_slide(
        prs,
        "전처리 방법론(P-track)",
        [
            (
                "P1 데이터 정합",
                [
                    "common-6 라벨 매핑으로 CREMA-D/RAVDESS 통합",
                    "manifest/split 자동 생성으로 재현성 확보",
                    "AV 중심, AO 보조 비교셋 병행",
                ],
            ),
            (
                "P2 표현 고도화",
                [
                    "cache_v2: log-mel/delta 통계 특징",
                    "cache_v3: raw 재추출 + pretrained video",
                    "cache_v5: HuBERT 오디오 표현 강화",
                ],
            ),
            (
                "P3 정합/품질",
                [
                    "main/cross 공통 규약으로 비교 가능성 확보",
                    "캐시 버전 전략으로 변경 영향 분리",
                    "미완료: ROI/품질플래그/CMVN/SpecAug",
                ],
            ),
            (
                "P4 연구적 의미",
                [
                    "입력 표현이 성능 병목임을 실증",
                    "모델 개선과 입력 개선의 기여 분해 가능",
                    "0.9 트랙 품질 인지형 전처리의 근거 확보",
                ],
            ),
        ],
    )

    add_2x2_blocks_slide(
        prs,
        "모델 설계/학습 방법론(A/T-track)",
        [
            (
                "A1 비교군 설계",
                [
                    "B0 -> ML(LR/SVM/RF) -> FP32 멀티태스크",
                    "모달(audio/video/fusion)별 동시 비교",
                    "단계별 기준선 유지로 delta 해석 가능",
                ],
            ),
            (
                "A2 FP32 구조",
                [
                    "dual-encoder + fusion(concat/gated)",
                    "emotion + arousal2/3 multi-head",
                    "결측 라벨(-1) 마스킹으로 샘플 손실 최소화",
                ],
            ),
            (
                "A3 학습전략",
                [
                    "CE/Focal + weighted sampler + label smoothing",
                    "cosine scheduler + grad clipping",
                    "gated/wide 확장으로 결합 용량 강화",
                ],
            ),
            (
                "T/Domain",
                [
                    "audio pretrained: wav2vec2 -> HuBERT",
                    "video pretrained embedding 결합",
                    "cross 적응: CORAL none vs coral 비교",
                ],
            ),
        ],
    )

    add_2x2_blocks_slide(
        prs,
        "비교/평가/통계 방법론",
        [
            (
                "E1 지표",
                [
                    "주지표: emotion macro-F1",
                    "보조: accuracy/OVR-AUC, arousal2/3 MAE",
                    "불균형에서 accuracy 편향 보정",
                ],
            ),
            (
                "E2 통제",
                [
                    "동일 분할·라벨·지표 유지",
                    "변수(전처리/모델/학습) 단계별 분리",
                    "main/cross 동시 관찰",
                ],
            ),
            (
                "E3 통계 검증",
                [
                    "phase2 pairwise bootstrap CI",
                    "phase3 vs phase2 delta CI",
                    "유의한 개선/비개선 구간 분리",
                ],
            ),
            (
                "E4 운영 재현성",
                [
                    "summary/predictions/progress 표준화",
                    "장시간 무출력 이슈 -> progress.json",
                    "스크립트 기반 리포트 자동 생성",
                ],
            ),
        ],
    )

    cache_chain_rows: List[List[str]] = []
    for name, info in [
        ("cache_v2", cache_v2),
        ("cache_v3", cache_v3),
        ("cache_v4(wav2vec2)", cache_v4),
        ("cache_v5_hubert", cache_v5),
    ]:
        cache_chain_rows.append(
            [
                name,
                str(info.get("audio_pretrained_backbone", "none")),
                f"{int(as_float(info.get('ok')))} / {int(as_float(info.get('total_rows')))}",
                str(info.get("device_resolved", info.get("device", "-"))),
                f"{as_float(info.get('elapsed_sec')):.1f}",
            ]
        )
    add_table_slide(
        prs,
        "특징 캐시 체인 상세(무엇을 어디에 적용했는가)",
        ["캐시", "오디오 백본", "성공/전체", "장치", "생성시간(s)"],
        cache_chain_rows,
        col_widths=[2.8, 2.4, 2.0, 2.2, 2.0],
        top=1.45,
        height=4.7,
    )

    add_table_slide(
        prs,
        "핵심 설정 Traceability 매트릭스",
        ["설정/방법", "코드/스크립트 적용 위치", "무엇이 바뀌었나", "결과(요약)"],
        [
            [
                "audio_pretrained_backbone=wav2vec2_base",
                "scripts/run_phase35_strong_v1.sh (STEP1), cache_v4",
                "오디오 임베딩을 통계 특징에 추가",
                f"v5 main {m['v5_main_f1']:.4f}",
            ],
            [
                "audio_pretrained_backbone=hubert_base",
                "scripts/run_phase35_next_v8_hubert_main.sh (STEP1), cache_v5_hubert",
                "오디오 표현력 상향",
                f"v8 단일 {m['v8_single_f1']:.4f}",
            ],
            [
                "fusion_type=gated + modality_dropout",
                "scripts/train_fp32_multitask.py (--fusion-type, --modality-dropout-p)",
                "모달 신뢰도 차이를 반영하는 결합",
                f"v7 {m['v7_main_f1']:.4f} -> v8 {m['v8_single_f1']:.4f}",
            ],
            [
                "label_smoothing + weighted_sampler",
                "scripts/train_fp32_multitask.py (--label-smoothing, --weighted-sampler)",
                "불균형/과신 문제 완화",
                f"v6 {m['v6_main_f1']:.4f}, v7 {m['v7_main_f1']:.4f}",
            ],
            [
                "domain_adapt=coral",
                "scripts/train_ml_baselines.py (--domain-adapt coral)",
                "도메인 공분산 정렬",
                f"cross {m['v8_cross_c2r_f1']:.4f}/{m['v8_cross_r2c_f1']:.4f}",
            ],
        ],
        col_widths=[2.9, 4.2, 2.8, 2.3],
        top=1.25,
        height=5.45,
    )

    add_table_slide(
        prs,
        "전처리/표현 개념 이론 + 실제 적용",
        ["개념", "이론/의도", "실제 적용(설정)", "상태"],
        [
            ["log-mel", "주파수 축을 지각 스케일로 압축", "cache_v2 audio_feature_v2", "적용"],
            ["delta feature", "시간 변화율로 억양 변화 포착", "cache_v2 d1 통계", "적용"],
            ["pretrained audio embedding", "대규모 음성표현 prior 활용", "wav2vec2/hubert/wavlm 옵션", "적용"],
            ["pretrained video embedding", "표정/장면 표현력 강화", "resnet18/resnet34/efficientnet_b0", "적용"],
            ["CMVN / SpecAugment", "도메인 편차/과적합 완화", "P-track 강한 버전 항목", "미완료"],
            ["ROI/품질플래그", "얼굴 검출 품질 기반 신뢰도 반영", "quality-aware fusion 입력", "미완료"],
        ],
        col_widths=[2.4, 3.2, 3.8, 1.6],
        top=1.2,
        height=5.6,
    )

    add_table_slide(
        prs,
        "모델/학습 개념 이론 + 실제 적용",
        ["개념", "핵심 이론", "적용 파라미터", "적용 위치"],
        [
            ["Cross-Entropy", "클래스 확률 최대화 기본 목적함수", "--emotion-loss ce", "train_fp32_multitask.py"],
            ["Focal Loss", "쉬운 샘플 가중치 감소, 어려운 샘플 집중", "--emotion-loss focal --focal-gamma", "train_fp32_multitask.py"],
            ["Weighted Sampler", "불균형 클래스 출현 빈도 보정", "--weighted-sampler", "train_fp32_multitask.py"],
            ["Label Smoothing", "과신(confidence overfit) 완화", "--label-smoothing", "train_fp32_multitask.py"],
            ["Gated Fusion", "모달별 신뢰도 가중 결합", "--fusion-type gated", "train_fp32_multitask.py"],
            ["Cosine LR", "학습률을 주기적으로 완만 감소", "CosineAnnealingLR", "train_fp32_multitask.py"],
        ],
        col_widths=[2.3, 3.4, 2.7, 3.7],
        top=1.2,
        height=5.6,
    )

    add_2x2_blocks_slide(
        prs,
        "CORAL 도메인 적응: 개념 -> 구현 -> 결과",
        [
            (
                "개념/이론",
                [
                    "source/target feature covariance 차이를 줄여 도메인 정렬",
                    "분류기 구조는 유지하고 입력 분포를 보정",
                    "데이터셋 스타일 차이에 강건한 baseline 목적",
                ],
            ),
            (
                "구현",
                [
                    "scripts/train_ml_baselines.py",
                    "--domain-adapt none vs coral 비교",
                    "logreg, linear_svm에서 동일 프로토콜 적용",
                ],
            ),
            (
                "실험 설계",
                [
                    "CREMA->RAVDESS, RAVDESS->CREMA 양방향",
                    "phase2/v3/v5/v8 단계별 반복 측정",
                    "bootstrap 보고서와 함께 해석",
                ],
            ),
            (
                "결과",
                [
                    f"C->R: {m['phase2_cross_c2r_f1']:.4f} -> {m['v8_cross_c2r_f1']:.4f}",
                    f"R->C: {m['phase2_cross_r2c_f1']:.4f} -> {m['v8_cross_r2c_f1']:.4f}",
                    "적응 미적용 대비 명확한 개선 확인",
                ],
            ),
        ],
    )

    v7_rows = [
        ["epochs", str(v7_wide_run.get("epochs", "-")), str(v8_tune4_run.get("epochs", "-")), "학습 길이"],
        ["lr", str(v7_wide_run.get("lr", "-")), str(v8_tune4_run.get("lr", "-")), "수렴 속도/안정성"],
        ["hidden_dim", str(v7_wide_run.get("hidden_dim", "-")), str(v8_tune4_run.get("hidden_dim", "-")), "표현 용량"],
        ["emb_dim", str(v7_wide_run.get("emb_dim", "-")), str(v8_tune4_run.get("emb_dim", "-")), "모달 임베딩 크기"],
        ["dropout", str(v7_wide_run.get("dropout", "-")), str(v8_tune4_run.get("dropout", "-")), "과적합 억제"],
        ["modality_dropout_p", str(v7_wide_run.get("modality_dropout_p", "-")), str(v8_tune4_run.get("modality_dropout_p", "-")), "모달 강건성"],
        ["label_smoothing", str(v7_wide_run.get("label_smoothing", "-")), str(v8_tune4_run.get("label_smoothing", "-")), "과신 완화"],
        ["weighted_sampler", str(v7_wide_run.get("weighted_sampler", "-")), str(v8_tune4_run.get("weighted_sampler", "-")), "불균형 대응"],
        ["device", str(v7_wide_run.get("device", v7_wide_run.get("device_resolved", "-"))), str(v8_tune4_run.get("device_resolved", v8_tune4_run.get("device", "-"))), "실행 환경"],
    ]
    add_table_slide(
        prs,
        "v7(wide) vs v8(tune4) 하이퍼파라미터 상세",
        ["설정", "v7", "v8", "의미"],
        v7_rows,
        col_widths=[2.5, 2.7, 2.7, 4.1],
        top=1.2,
        height=5.6,
    )

    add_table_slide(
        prs,
        "실험 파이프라인 상세(무엇을 어디에 어떻게)",
        ["단계", "실행 스크립트/코드", "입력/설정", "출력/결과"],
        [
            [
                "데이터 정합",
                "scripts/prepare_research_data.py",
                "CREMA-D + RAVDESS, common-6",
                "manifest/splits/summary.json",
            ],
            [
                "특징 생성",
                "scripts/prepare_advanced_features.py",
                "cache_v2/v3/v4/v5, pretrained backbone 옵션",
                "derived/features/cache_*/summary.json",
            ],
            [
                "ML baseline",
                "scripts/train_ml_baselines.py",
                "logreg/rbf_svm/linear_svm, domain_adapt",
                "phase2/phase35 *_metrics.csv",
            ],
            [
                "FP32 multitask",
                "scripts/train_fp32_multitask.py",
                "fusion/loss/sampler/smoothing/hparams",
                "summary.json + predictions.csv",
            ],
            [
                "분석/리포트",
                "scripts/analyze_phase*.py",
                "bootstrap + delta 비교",
                "phase*_results.md / *_metrics.csv",
            ],
        ],
        col_widths=[1.9, 3.4, 3.6, 3.3],
        top=1.25,
        height=5.5,
    )

    # Section III
    add_section_divider(prs, "III. 단계별 실행 및 검증", "각 단계를 가설-방법-비교-인사이트 템플릿으로 정리")

    add_table_slide(
        prs,
        "단계별 실행 로드맵(요약)",
        ["단계", "핵심 질문", "핵심 방법", "대표 성과"],
        [
            ["Phase-1", "기준선은 어디인가?", "majority baseline", "F1 약 0.05"],
            ["Phase-2", "fusion 이득이 있는가?", "ML(audio/video/fusion)", "main fusion 0.3950"],
            ["Phase-3", "FP32 multitask 효과는?", "dual-encoder", "main 하락/cross 일부 개선"],
            ["v2~v3", "표현 고도화 효과는?", "cache 개선 + pretrained", "main 0.4807"],
            ["v4~v5", "적응으로 cross 회복 가능한가?", "CORAL + strong-v1", "cross 0.3025/0.2724"],
            ["v6~v8", "0.7 달성 가능한가?", "recipe+gated+HuBERT", "단일 0.6992/앙상블 0.7099"],
        ],
        col_widths=[1.9, 3.7, 3.6, 3.2],
    )

    add_phase_card(
        prs,
        "Phase-1(B0): 기준선 확립",
        "최빈 기준선을 통해 개선폭을 얼마나 신뢰성 있게 해석할 수 있는가?",
        ["majority baseline 구성", "macro-F1 중심 비교축 고정"],
        ["B0 F1 약 0.05", "후속 단계 delta 비교 anchor 확보"],
        ["성능 목표보다 평가 체계의 기준점 확보가 목적", "이후 개선폭 해석 신뢰성 확보"],
    )

    p2_main = {r["mode"]: r for r in m["phase2_main_modes"]}
    p2_c2r = {r["mode"]: r for r in m["phase2_cross_modes"]["cross_crema_to_ravdess"]}
    p2_r2c = {r["mode"]: r for r in m["phase2_cross_modes"]["cross_ravdess_to_crema"]}

    add_phase_card(
        prs,
        "Phase-2: 멀티모달 기본 성능 검증",
        "고전 모델에서도 fusion이 일관된 이득을 보이는가?",
        ["LR/SVM/RF를 audio/video/fusion으로 병렬 비교", "main/cross 양방향 동시 측정", "동일 프로토콜로 통제"],
        [
            f"main F1: audio {p2_main['audio']['f1']:.4f}, video {p2_main['video']['f1']:.4f}, fusion {p2_main['fusion']['f1']:.4f}",
            f"cross C->R fusion {p2_c2r['fusion']['f1']:.4f}, R->C fusion {p2_r2c['fusion']['f1']:.4f}",
        ],
        ["main에서는 fusion 이득이 확인", "cross는 방향 비대칭으로 도메인 갭이 핵심 병목"],
    )

    phase2_boot_rows = []
    for r in m["phase2_pairwise_main"]:
        phase2_boot_rows.append(
            [
                f"{r['lhs']} - {r['rhs']}",
                f"{r['delta']:+.4f}",
                format_ci(r["ci_lo"], r["ci_hi"]),
                "유의" if (r["ci_lo"] > 0 or r["ci_hi"] < 0) else "불명확",
            ]
        )
    add_table_slide(
        prs,
        "Phase-2 pairwise bootstrap(main)",
        ["비교", "delta(F1)", "95% CI", "해석"],
        phase2_boot_rows,
        col_widths=[3.6, 1.8, 3.0, 1.8],
        top=1.4,
        height=4.9,
    )

    add_phase_card(
        prs,
        "Phase-3(FP32): 멀티태스크 딥러닝 검증",
        "FP32 구조가 main/cross를 동시에 개선하는가?",
        ["dual-encoder + multitask", "결측 라벨 마스킹 처리", "학습 안정화 recipe 반영"],
        [
            f"main F1 {m['phase3_main_f1']:.4f}",
            f"cross C->R {m['phase3_cross']['cross_crema_to_ravdess']:.4f}, R->C {m['phase3_cross']['cross_ravdess_to_crema']:.4f}",
        ],
        ["구조만으로 main 우세가 자동 보장되지 않음", "표현/학습/적응을 함께 올려야 함"],
    )

    boot = m["phase3_bootstrap"]
    p3_rows = []
    for run_key, run_label in [
        ("main", "main"),
        ("cross_crema_to_ravdess", "cross C->R"),
        ("cross_ravdess_to_crema", "cross R->C"),
    ]:
        r = boot.get(run_key, {})
        delta = as_float(r.get("delta"))
        ci_lo = as_float(r.get("ci_lo"))
        ci_hi = as_float(r.get("ci_hi"))
        p3_rows.append(
            [
                run_label,
                f"{delta:+.4f}",
                format_ci(ci_lo, ci_hi),
                str(int(as_float(r.get("n")))),
                "유의" if (ci_lo > 0 or ci_hi < 0) else "불명확",
            ]
        )
    add_table_slide(
        prs,
        "Phase-3 vs Phase-2(fusion) bootstrap delta",
        ["run", "delta(F1)", "95% CI", "n", "판정"],
        p3_rows,
        col_widths=[2.6, 2.0, 3.2, 1.2, 1.5],
        top=1.6,
        height=4.3,
    )

    add_phase_card(
        prs,
        "Phase-3.5(v2~v3): 표현 고도화",
        "입력 표현 강화가 성능 병목을 해소하는가?",
        ["cache_v2(log-mel/delta)", "cache_v3(pretrained video)", "ML/FP32 병행 비교"],
        [f"v2 main 최고 {m['v2_main_f1']:.4f}", f"v3 main 최고 {m['v3_main_f1']:.4f}", "cross baseline은 일부 붕괴"],
        ["표현 개선은 main 성능에 직접 기여", "적응이 없으면 cross 일반화는 취약"],
    )

    add_phase_card(
        prs,
        "Phase-3.5(v4~v5): 도메인 적응 + strong-v1",
        "CORAL이 cross 붕괴를 회복하는가?",
        ["same classifier(logreg)에서 none vs coral 비교", "audio pretrained 강화 결합", "main/cross 동시 모니터링"],
        [
            f"v5 main {m['v5_main_f1']:.4f}",
            f"v5+CORAL C->R {m['v5_cross_c2r_f1']:.4f}",
            f"v5+CORAL R->C {m['v5_cross_r2c_f1']:.4f}",
        ],
        ["CORAL은 cross 회복에 실효적", "표현 강화 + 도메인 정렬의 결합이 핵심"],
    )

    ca = m["cross_ablation"]
    add_table_slide(
        prs,
        "Cross-domain ablation (none vs CORAL)",
        ["방향", "phase2 fusion", "v3 baseline", "v5 baseline", "v5+CORAL", "v8+CORAL"],
        [
            [
                "CREMA->RAVDESS",
                f"{ca['crema_to_ravdess']['phase2_fusion']:.4f}",
                f"{ca['crema_to_ravdess']['v3_baseline']:.4f}",
                f"{ca['crema_to_ravdess']['v5_baseline']:.4f}",
                f"{ca['crema_to_ravdess']['v5_coral']:.4f}",
                f"{ca['crema_to_ravdess']['v8_coral']:.4f}",
            ],
            [
                "RAVDESS->CREMA",
                f"{ca['ravdess_to_crema']['phase2_fusion']:.4f}",
                f"{ca['ravdess_to_crema']['v3_baseline']:.4f}",
                f"{ca['ravdess_to_crema']['v5_baseline']:.4f}",
                f"{ca['ravdess_to_crema']['v5_coral']:.4f}",
                f"{ca['ravdess_to_crema']['v8_coral']:.4f}",
            ],
        ],
        col_widths=[2.8, 1.8, 1.8, 1.8, 1.8, 1.8],
        top=1.8,
        height=3.7,
    )

    add_phase_card(
        prs,
        "Phase-3.5(v6~v8): 0.7 구간 진입",
        "학습전략 고도화 + HuBERT로 0.7 목표에 도달하는가?",
        [f"v6 {m['v6_main_f1']:.4f}", f"v7 {m['v7_main_f1']:.4f}", "v8 HuBERT tune + ensemble"],
        [
            f"v8 단일 {m['v8_single_f1']:.4f}",
            f"v8 앙상블 {m['v8_ensemble_f1']:.4f}",
            f"cross 최고 {m['v8_cross_c2r_f1']:.4f}/{m['v8_cross_r2c_f1']:.4f}",
        ],
        ["0.7은 앙상블 기준 달성", "단일모델은 임계점 근접, 품질/튜닝 고도화 시 상회 가능"],
    )

    add_phase2_modality_chart(prs, m)
    add_main_trend_chart(prs, m)
    add_cross_chart(prs, m)

    stage_rows = []
    prev = None
    for row in m["main_stage_metrics"]:
        delta = "-" if prev is None else f"{row['f1'] - prev:+.4f}"
        stage_rows.append([
            row["stage"],
            f"{row['f1']:.4f}",
            delta,
            f"{row['acc']:.4f}",
            f"{row['ar2']:.4f}",
            f"{row['ar3']:.4f}",
        ])
        prev = row["f1"]
    add_table_slide(
        prs,
        "main 대표 실험의 F1/MAE 추세",
        ["stage", "macro-F1", "delta(prev)", "accuracy", "arousal2 MAE", "arousal3 MAE"],
        stage_rows,
        col_widths=[2.4, 1.6, 1.8, 1.6, 2.0, 2.0],
        top=1.25,
        height=5.4,
    )

    top_rows = []
    for row in m["top_v8"]:
        top_rows.append(
            [
                row.get("name", ""),
                f"{as_float(row.get('emotion_macro_f1')):.4f}",
                f"{as_float(row.get('emotion_acc')):.4f}",
                f"{as_float(row.get('arousal2_mae')):.4f}",
                row.get("device", ""),
            ]
        )
    add_table_slide(
        prs,
        "v8 후보군 상세 비교(상위 8)",
        ["실험명", "F1", "Acc", "A2 MAE", "장치"],
        top_rows,
        col_widths=[6.7, 1.3, 1.3, 1.4, 1.5],
        top=1.2,
        height=5.6,
    )

    add_table_slide(
        prs,
        "핵심 장애와 해결(실제 운영 이슈)",
        ["이슈", "원인", "해결", "효과"],
        [
            ["WSL 의존성 누락", "numpy/pip 부재", "ensurepip/get-pip + 자동설치", "실험 재실행 안정화"],
            ["결측 라벨 처리", "학습 샘플 과도 누락", "-1 마스킹 방식", "fold 샘플 수 정상화"],
            ["GPU 미활용", "CPU torch/device 해석 문제", "cu126 전환 + 코드 수정", "CUDA 학습 정상화"],
            ["진행률 불투명", "장시간 무출력", "progress.json + 로그 단계화", "운영 가시성 개선"],
        ],
        col_widths=[2.6, 3.0, 3.3, 3.0],
        top=1.35,
        height=5.1,
    )

    add_table_slide(
        prs,
        "재현성 산출물 경로(핵심)",
        ["범주", "대표 파일/경로", "역할"],
        [
            ["설계 문서", "deep-research-report-advancement (3).md", "고도화 트랙/목표 정의"],
            ["진행 문서", "derived/reports/project_progress_until_2026-02-20.md", "이슈/타임라인/상태"],
            ["종합 보고서", "derived/reports/full_research_process_and_results_until_2026-02-24_ko.md", "방법론+인사이트 종합"],
            ["main/cross 지표", "derived/reports/phase2_global_metrics.csv 등", "정량 비교 근거"],
            ["최종 후보", "derived/reports/phase35_next_v8_metrics.csv", "0.7 근접/달성 근거"],
            ["도메인 적응", "derived/reports/phase35_cross_domain_adapt_metrics.csv", "cross 개선 근거"],
        ],
        col_widths=[1.8, 6.4, 3.8],
        top=1.25,
        height=5.5,
    )

    add_table_slide(
        prs,
        "설계 대비 구현 현황(P/A/T/Domain)",
        ["트랙", "설계 요구", "현재 상태", "해석"],
        [
            ["P-track", "ROI/품질플래그 + CMVN/SpecAug", "부분 적용", "mel/delta/임베딩 반영, quality-aware 전처리 미완료"],
            ["A-track", "boosting + 고급 fusion", "부분 적용", "LR/SVM/RF/gated 반영, XGB/LGBM/attention-lite 미완료"],
            ["T-track", "pretrained fine-tuning", "부분 적용", "HuBERT 기반 강화 성공, backbone unfreeze 제한적"],
            ["Domain", "cross 일반화 강화", "적용", "CORAL로 양방향 cross 일관 개선 확인"],
        ],
        col_widths=[1.6, 3.4, 1.8, 5.6],
    )

    # Section IV
    add_section_divider(prs, "IV. 결론과 의사결정", "0.7 달성 해석, 0.9 전환 조건, 실행 리스크 관리")

    add_table_slide(
        prs,
        "목표 트랙 분리: 0.7 vs 0.9",
        ["구분", "현재 최고", "목표", "격차", "판정"],
        [
            ["0.7 단일", f"{m['v8_single_f1']:.4f}", "0.7000", f"{0.7-float(m['v8_single_f1']):+.4f}", "근접"],
            ["0.7 앙상블", f"{m['v8_ensemble_f1']:.4f}", "0.7000", f"{0.7-float(m['v8_ensemble_f1']):+.4f}", "달성"],
            ["0.9 단일", f"{m['v8_single_f1']:.4f}", "0.9000", f"{gap_single:+.4f}", "미달"],
            ["0.9 앙상블", f"{m['v8_ensemble_f1']:.4f}", "0.9200", f"{gap_ens:+.4f}", "미달"],
        ],
        col_widths=[3.0, 2.1, 1.7, 2.0, 1.9],
        top=1.55,
        height=4.5,
    )

    add_2x2_blocks_slide(
        prs,
        "0.9 목표 전환 로드맵(설계 기반)",
        [
            (
                "R1 전처리 고도화",
                [
                    "ROI crop + 품질플래그(blur/brightness/success)",
                    "CMVN/SpecAug/loudness norm full pipeline",
                    "quality-aware feature fusion 입력 확장",
                ],
            ),
            (
                "R2 모델 구조 고도화",
                [
                    "cross-attention-lite 결합 실험",
                    "boosting(XGB/LGBM) classical track 추가",
                    "가벼운 구조와 성능의 Pareto 탐색",
                ],
            ),
            (
                "R3 사전학습 파인튜닝",
                [
                    "frozen embedding -> partial unfreeze",
                    "오디오/비디오 상위 레이어 감정특화 미세조정",
                    "seed 반복으로 분산 관리",
                ],
            ),
            (
                "R4 도메인 적응 강화",
                [
                    "CORAL beyond: feature-level adaptation",
                    "consistency/self-training 경로 도입",
                    "main + cross + 분산 동시 만족 기준 적용",
                ],
            ),
        ],
    )

    add_table_slide(
        prs,
        "리스크와 대응 전략",
        ["리스크", "영향", "완화 전략", "판정 기준"],
        [
            ["main overfit", "cross 하락", "cross 조기 경보 + 적응 병행", "main 상승과 cross 동시 유지"],
            ["복잡도 증가", "온디바이스 이행 지연", "구조-성능 Pareto 추적", "latency/FPS 목표 유지"],
            ["실험 분산 과대", "재현성 저하", "seed 반복 + CI 보고", "CI 폭/평균 동시 관리"],
            ["장시간 실행 중단", "실험 누락", "progress 모니터링 + 완주 정책", "실험군 누락 0건"],
        ],
        col_widths=[2.3, 2.2, 4.2, 3.5],
        top=1.35,
        height=5.2,
    )

    p95 = as_float(runtime.get("latency_ms", {}).get("p95"), default=-1)
    fps = as_float(runtime.get("fps_equiv"), default=-1)
    runtime_line = "phase2 runtime 벤치 값 미확인"
    if p95 >= 0 and fps >= 0:
        runtime_line = f"phase2 fold0 기준: p95 {p95:.2f} ms, FPS {fps:.2f}"

    add_bullet_slide(
        prs,
        "최종 정리",
        "현재 상태와 다음 의사결정",
        [
            "본 발표자료는 설계서/진행문서/종합보고서/실제 결과 파일을 교차해 재구성했습니다.",
            "핵심 결론: 0.7 트랙 달성, 0.9 트랙은 구조 전환형 고도화가 필수입니다.",
            runtime_line,
            "사용자 요청에 따라 신규 실험은 보류 상태이며, 보고/발표 자료 고도화를 우선 완료했습니다.",
            "질의응답",
        ],
        size=18,
    )

    prs.save(OUT_PPTX)


# ========================
# PDF build
# ========================


def pick_korean_fonts() -> Tuple[str | None, str | None]:
    regular_candidates = [
        "/mnt/c/Windows/Fonts/malgun.ttf",
        "/mnt/c/Windows/Fonts/HANDotum.ttf",
        "/mnt/c/Windows/Fonts/HanSantteutDotum-Regular.ttf",
    ]
    bold_candidates = [
        "/mnt/c/Windows/Fonts/malgunbd.ttf",
        "/mnt/c/Windows/Fonts/HANDotumB.ttf",
        "/mnt/c/Windows/Fonts/HanSantteutDotum-Bold.ttf",
    ]

    regular = next((p for p in regular_candidates if Path(p).exists()), None)
    bold = next((p for p in bold_candidates if Path(p).exists()), None)
    return regular, bold


def wrap_line(c: canvas.Canvas, text: str, font_name: str, font_size: int, max_width: float) -> List[str]:
    words = text.split(" ")
    if not words:
        return [""]

    lines: List[str] = []
    cur = words[0]
    for w in words[1:]:
        trial = cur + " " + w
        if c.stringWidth(trial, font_name, font_size) <= max_width:
            cur = trial
        else:
            lines.append(cur)
            cur = w
    lines.append(cur)
    return lines


def draw_block(c: canvas.Canvas, x: float, y: float, title: str, lines: List[str], font_r: str, font_b: str) -> float:
    c.setFont(font_b, 13)
    c.drawString(x, y, title)
    y -= 7 * mm

    c.setFont(font_r, 10)
    max_w = 170 * mm
    for line in lines:
        wrapped = wrap_line(c, line, font_r, 10, max_w)
        for w in wrapped:
            c.drawString(x + 2 * mm, y, w)
            y -= 5.2 * mm
    y -= 2 * mm
    return y


def build_pdf(m: Dict[str, object]) -> None:
    regular_path, bold_path = pick_korean_fonts()
    font_r = "Helvetica"
    font_b = "Helvetica-Bold"

    if regular_path and bold_path:
        pdfmetrics.registerFont(TTFont("KoRegular", regular_path))
        pdfmetrics.registerFont(TTFont("KoBold", bold_path))
        font_r = "KoRegular"
        font_b = "KoBold"

    c = canvas.Canvas(str(OUT_PDF), pagesize=A4)
    _, h = A4

    x = 18 * mm
    y = h - 18 * mm

    c.setFont(font_b, 18)
    c.drawString(x, y, "연구 진행 종합 요약 (상세 보강판)")
    y -= 8 * mm
    c.setFont(font_r, 10)
    c.drawString(x, y, "기준일: 2026-02-24 | 범위: 데이터 준비 ~ Phase-3.5(v8) + 설계 대비 구현 점검")
    y -= 10 * mm

    gap_single = 0.9 - float(m["v8_single_f1"])
    gap_ens = 0.92 - float(m["v8_ensemble_f1"])
    sections = [
        (
            "1) 핵심 수치",
            [
                f"- Main 단일모델 최고: {m['v8_single_f1']:.4f}",
                f"- Main 앙상블 최고: {m['v8_ensemble_f1']:.4f}",
                f"- Cross 최고(C->R / R->C): {m['v8_cross_c2r_f1']:.4f} / {m['v8_cross_r2c_f1']:.4f}",
            ],
        ),
        (
            "2) 방법론 요약",
            [
                "- 전처리: cache_v1~v5_hubert 버전 전략으로 표현 고도화 효과 분리",
                "- 모델: B0/ML/FP32 멀티태스크 계층 비교",
                "- 평가: main/cross 분리 + macro-F1 중심 + bootstrap CI",
                "- 운영: progress.json 기반 장시간 실험 가시화",
            ],
        ),
        (
            "3) 설정/이론 연결",
            [
                "- wav2vec2/hubert, gated fusion, label smoothing, weighted sampler 적용 위치를 코드 단위로 추적",
                "- CORAL의 도메인 공분산 정렬 개념과 cross 개선 결과를 연결",
                "- v7-v8 하이퍼파라미터 차이를 실제 summary.json 기반으로 비교",
            ],
        ),
        (
            "4) 단계별 결론",
            [
                "- Phase-2: main fusion 이득 확인, cross 비대칭 확인",
                "- Phase-3: 구조 단독 개선의 한계, 표현/학습 결합 필요",
                "- v4/v5: CORAL로 cross 회복",
                "- v8: HuBERT + gated + ensemble로 0.7 달성 구간 진입",
            ],
        ),
        (
            "5) 설계 대비 구현 상태",
            [
                "- 반영: pretrained 임베딩, gated fusion, CORAL, GPU 경로",
                "- 미완료: ROI/품질플래그, CMVN/SpecAug, attention-lite, unfreeze fine-tuning",
                "- 결론: 강한 버전은 부분 구현 상태",
            ],
        ),
        (
            "6) 의사결정",
            [
                f"- 0.7 트랙: {m['v8_ensemble_f1']:.4f} (달성)",
                f"- 0.9 단일모델 격차: {gap_single:+.4f}",
                f"- 0.9 앙상블(0.92) 격차: {gap_ens:+.4f}",
                "- 권고: 전처리/모델/적응의 구조 전환형 고도화 후 실험 재개",
            ],
        ),
        (
            "7) 현재 상태",
            [
                "- 사용자 요청으로 신규 실험은 보류",
                "- 문서/발표자료를 방법론+정량근거 중심으로 보강 완료",
                "- 실험 재개 시 0.7/0.9 트랙 분리 운영 권장",
            ],
        ),
    ]

    for title, lines in sections:
        y = draw_block(c, x, y, title, lines, font_r, font_b)
        if y < 35 * mm:
            c.showPage()
            y = h - 18 * mm
            c.setFont(font_r, 10)

    c.save()


# ========================
# Entrypoint
# ========================


def main() -> None:
    REPORTS.mkdir(parents=True, exist_ok=True)
    metrics = load_metrics()
    build_pptx(metrics)
    build_pdf(metrics)
    print(
        json.dumps(
            {
                "pptx": str(OUT_PPTX),
                "pdf": str(OUT_PDF),
                "phase2_main_f1": metrics["phase2_main_f1"],
                "phase3_main_f1": metrics["phase3_main_f1"],
                "v8_single_f1": metrics["v8_single_f1"],
                "v8_ensemble_f1": metrics["v8_ensemble_f1"],
            },
            ensure_ascii=True,
            indent=2,
        )
    )


if __name__ == "__main__":
    main()
